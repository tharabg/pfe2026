"""
Génère la présentation PFE TOPNET OCR en PowerPoint professionnel
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Couleurs : blanc + bleu moyen ────────────────────────────────────────────
BLUE_DARK   = RGBColor(0x1A, 0x56, 0x9E)   # bleu moyen principal  #1A569E
BLUE_MED    = RGBColor(0x2E, 0x78, 0xC5)   # bleu clair secondaire #2E78C5
BLUE_LIGHT  = RGBColor(0xDB, 0xEA, 0xFE)   # fond bleu très clair  #DBEAFE
ORANGE      = RGBColor(0xff, 0x95, 0x00)   # accent orange         #ff9500
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT  = RGBColor(0xF5, 0xF8, 0xFF)   # fond slide principal
GRAY_TEXT   = RGBColor(0x33, 0x4E, 0x7A)   # texte foncé sur blanc
GREEN_OK    = RGBColor(0x16, 0xA3, 0x4A)
RED_KO      = RGBColor(0xDC, 0x26, 0x26)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def rect(slide, x, y, w, h, fill=None, line=None, radius=False):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def txbox(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb

def multiline_box(slide, lines, x, y, w, h, size=14, color=WHITE, bold_first=False):
    """lines = list of (text, bold, color_override)"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, col = item, False, color
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            col  = item[2] if len(item) > 2 else color
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = col
        run.font.name = "Calibri"
    return tb

def header_bar(slide, title, subtitle=None):
    """Bande bleue en haut avec titre orange"""
    rect(slide, 0, 0, 13.33, 1.3, fill=BLUE_DARK)
    rect(slide, 0, 1.3, 13.33, 0.06, fill=ORANGE)
    txbox(slide, title, 0.4, 0.15, 11, 0.7, size=28, bold=True,
          color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txbox(slide, subtitle, 0.4, 0.75, 11, 0.5, size=14,
              color=ORANGE, align=PP_ALIGN.LEFT)

def footer(slide, text="Ben Ghorbel Thara — PFE TOPNET 2026"):
    rect(slide, 0, 7.1, 13.33, 0.4, fill=BLUE_DARK)
    txbox(slide, text, 0.3, 7.12, 8, 0.3, size=10, color=BLUE_LIGHT)
    txbox(slide, "TOPNET OCR", 10.5, 7.12, 2.5, 0.3, size=10,
          bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)

def card(slide, x, y, w, h, title, lines, title_color=ORANGE, bg=BLUE_MED):
    rect(slide, x, y, w, h, fill=bg)
    rect(slide, x, y, w, 0.45, fill=BLUE_DARK)
    txbox(slide, title, x+0.15, y+0.05, w-0.3, 0.38,
          size=13, bold=True, color=title_color)
    items = [(f"  {l}", False, WHITE) if not isinstance(l, tuple) else l for l in lines]
    multiline_box(slide, items, x+0.1, y+0.52, w-0.2, h-0.6, size=11.5, color=WHITE)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITRE
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
rect(sl, 0, 0, 13.33, 3.2, fill=BLUE_DARK)
rect(sl, 0, 3.2, 13.33, 0.07, fill=ORANGE)

# Bande gauche accent
rect(sl, 0, 0, 0.35, 7.5, fill=ORANGE)

txbox(sl, "TOPNET OCR", 0.6, 0.25, 9, 0.85, size=38, bold=True, color=WHITE)
txbox(sl, "Système Intelligent d'OCR", 0.6, 1.05, 12, 0.65,
      size=26, bold=True, color=WHITE)
txbox(sl, "et de Détection de Fraude Documentaire", 0.6, 1.65, 12, 0.55,
      size=20, bold=False, color=BLUE_LIGHT)
rect(sl, 0.6, 2.38, 4.5, 0.05, fill=ORANGE)

multiline_box(sl, [
    ("Stagiaire : Ben Ghorbel Thara", True,  GRAY_TEXT),
    ("Établissement : FST — Université de Tunis El Manar", False, GRAY_TEXT),
    ("Entreprise : TOPNET Tunisia  |  Stage PFE 6 mois  |  2026", False, GRAY_TEXT),
], 0.6, 3.55, 8.2, 1.4, size=14)

# Badge métriques
rect(sl, 9.2, 3.45, 3.8, 2.85, fill=BLUE_LIGHT)
rect(sl, 9.2, 3.45, 3.8, 0.42, fill=BLUE_DARK)
txbox(sl, "Métriques validées", 9.3, 3.48, 3.6, 0.35, size=12, bold=True, color=WHITE)
multiline_box(sl, [
    ("✅  Accuracy       85.7%", False, GRAY_TEXT),
    ("✅  Precision     100.0%", False, GRAY_TEXT),
    ("✅  F1-score       85.7%", False, GRAY_TEXT),
    ("✅  Faux positifs    0%",  True,  GREEN_OK),
], 9.3, 3.97, 3.6, 2.2, size=13)

footer(sl)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PLAN
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
rect(sl, 0, 0, 13.33, 7.5, fill=GRAY_LIGHT)
header_bar(sl, "Plan de la présentation")

items = [
    ("01", "Présentation de l'entreprise TOPNET"),
    ("02", "Contexte & Problématique"),
    ("03", "Solution proposée & Valeur ajoutée"),
    ("04", "Méthodologie"),
    ("05", "Besoins fonctionnels & non fonctionnels"),
    ("06", "Choix technologiques"),
    ("07", "Architecture de la solution"),
    ("08", "Avancement du projet"),
    ("09", "Conclusion & Perspectives"),
]
cols = [(0.4, items[:5]), (6.9, items[5:])]
for x_start, group in cols:
    for i, (num, label) in enumerate(group):
        y = 1.55 + i * 1.0
        rect(sl, x_start, y, 0.55, 0.7, fill=ORANGE)
        txbox(sl, num, x_start, y+0.1, 0.55, 0.5, size=16,
              bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, x_start+0.55, y, 5.9, 0.7, fill=WHITE)
        txbox(sl, label, x_start+0.7, y+0.15, 5.6, 0.45,
              size=14, color=BLUE_DARK)
footer(sl)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — TOPNET
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
rect(sl, 0, 0, 13.33, 7.5, fill=GRAY_LIGHT)
header_bar(sl, "Présentation de l'entreprise", "TOPNET Tunisia — Premier FAI privé")

# Carte gauche
rect(sl, 0.4, 1.55, 5.8, 5.25, fill=BLUE_LIGHT)
rect(sl, 0.4, 1.55, 5.8, 0.5, fill=BLUE_DARK)
txbox(sl, "TOPNET en chiffres", 0.55, 1.58, 5.5, 0.42,
      size=14, bold=True, color=WHITE)
multiline_box(sl, [
    ("🏢  Fondée en 1999 — 25+ ans d'activité", False, GRAY_TEXT),
    ("📍  Siège : Tunis | Présence nationale", False, GRAY_TEXT),
    ("🌐  Premier FAI privé de Tunisie", True,  BLUE_DARK),
    ("📶  Services : ADSL · FIBRE · 4G/5G · BOX", False, GRAY_TEXT),
    ("👥  Milliers de clients abonnés", False, GRAY_TEXT),
    ("📄  Docs traités/jour : CIN, Passeports,", False, GRAY_TEXT),
    ("     Factures, Contrats d'abonnement", False, GRAY_TEXT),
], 0.55, 2.15, 5.5, 4.4, size=13)

# Carte droite — problème actuel
rect(sl, 6.7, 1.55, 6.2, 5.25, fill=WHITE)
rect(sl, 6.7, 1.55, 6.2, 0.5, fill=BLUE_DARK)
txbox(sl, "Situation actuelle — Traitement documents",
      6.85, 1.58, 6.0, 0.42, size=13, bold=True, color=WHITE)
multiline_box(sl, [
    ("❌  Traitement 100% manuel", True,  RED_KO),
    ("     → 5 à 10 min par document", False, GRAY_TEXT),
    ("", False, WHITE),
    ("❌  Aucune détection de fraude auto.", True, RED_KO),
    ("     → Documents falsifiés non détectés", False, GRAY_TEXT),
    ("", False, WHITE),
    ("❌  Zéro traçabilité centralisée", True, RED_KO),
    ("     → Impossible d'analyser les données", False, GRAY_TEXT),
    ("", False, WHITE),
    ("💡  Besoin : automatiser, sécuriser,", True, BLUE_DARK),
    ("     analyser les données documents", True, BLUE_DARK),
], 6.85, 2.15, 5.9, 4.4, size=12.5)
footer(sl)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SOLUTION PROPOSÉE
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
rect(sl, 0, 0, 13.33, 7.5, fill=GRAY_LIGHT)
header_bar(sl, "Solution proposée", "TOPNET OCR — Pipeline IA local 7 étapes")

# 4 cartes fonctionnalités
cards_data = [
    ("🔍  Détection Fraude",    ["ELA (analyse pixels)", "VLM forensique Qwen2.5vl", "Score combiné par type doc"]),
    ("🏷️  Classification Auto", ["CLIP + VLM + Règles OCR", "4 types reconnus", "Fusion pondérée triple couche"]),
    ("📋  Extraction OCR",      ["PaddleOCR + EasyOCR", "Champs bilingues AR/FR", "Validation métier TOPNET"]),
    ("🌐  Interface Web",       ["Angular 21 — 2 rôles", "Upload · Résultats · Historique", "Design TOPNET professionnel"]),
]
for i, (title, pts) in enumerate(cards_data):
    x = 0.35 + i * 3.17
    rect(sl, x, 1.55, 3.0, 3.2, fill=BLUE_LIGHT)
    rect(sl, x, 1.55, 3.0, 0.55, fill=ORANGE)
    txbox(sl, title, x+0.1, 1.58, 2.8, 0.48, size=13, bold=True, color=BLUE_DARK)
    multiline_box(sl, [("• "+p, False, WHITE) for p in pts],
                  x+0.1, 2.2, 2.8, 2.4, size=12)

# Bande métriques
rect(sl, 0.35, 4.95, 12.63, 1.55, fill=BLUE_LIGHT)
rect(sl, 0.35, 4.95, 12.63, 0.45, fill=BLUE_DARK)
txbox(sl, "Métriques validées sur jeu de test (7 documents)", 0.55, 4.97,
      12, 0.4, size=13, bold=True, color=ORANGE)
metrics = [
    ("Accuracy\n85.7%", GREEN_OK), ("Precision\n100%", GREEN_OK),
    ("Recall\n75%", ORANGE),       ("F1-score\n85.7%", GREEN_OK),
    ("Faux positifs\n0%", GREEN_OK),("Temps moy.\n121s (CPU)", ORANGE),
]
for i, (txt, col) in enumerate(metrics):
    x = 0.6 + i * 2.05
    txbox(sl, txt, x, 5.45, 1.9, 0.95, size=12, bold=True,
          color=col, align=PP_ALIGN.CENTER)
footer(sl)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CHOIX TECHNOLOGIQUES
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
rect(sl, 0, 0, 13.33, 7.5, fill=GRAY_LIGHT)
header_bar(sl, "Choix technologiques", "Stack 100% local — Justifié")

techs = [
    ("FastAPI", "Backend Python", "Async, doc auto Swagger, performant"),
    ("Qwen2.5vl:7b", "VLM Vision", "Local, arabe+français, open-source"),
    ("CLIP ViT-Large", "Classification", "Zero-shot, 0 réentraînement"),
    ("PaddleOCR", "OCR arabe/fr", "Meilleur modèle open-source AR"),
    ("ELA + VLM", "Détection fraude", "Analyse pixel + forensique visuel"),
    ("Angular 21", "Frontend", "Enterprise-grade, maintenable"),
    ("PostgreSQL", "Base de données", "Prêt ETL + dashboards Power BI"),
    ("Ollama", "Runtime VLM", "Local, GPU/CPU, modèles interchangeables"),
]
for i, (name, cat, desc) in enumerate(techs):
    row, col = divmod(i, 4)
    x = 0.35 + col * 3.17
    y = 1.55 + row * 2.45
    rect(sl, x, y, 3.0, 2.15, fill=WHITE)
    rect(sl, x, y, 3.0, 0.38, fill=BLUE_DARK)
    txbox(sl, name, x+0.1, y+0.04, 2.8, 0.32, size=13, bold=True, color=ORANGE)
    txbox(sl, cat,  x+0.1, y+0.42, 2.8, 0.32, size=11, bold=True, color=BLUE_DARK)
    txbox(sl, desc, x+0.1, y+0.78, 2.8, 0.9,  size=10.5, color=GRAY_TEXT)
footer(sl)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ARCHITECTURE PIPELINE
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
rect(sl, 0, 0, 13.33, 7.5, fill=GRAY_LIGHT)
header_bar(sl, "Architecture — Pipeline IA 7 étapes")

steps = [
    ("01", "Détection Fraude",    "ELA + VLM forensique",       ORANGE),
    ("02", "Classification",      "CLIP + VLM + Règles OCR",    RGBColor(0x38,0xBD,0xF8)),
    ("03", "OCR Multi-moteur",    "PaddleOCR + EasyOCR",        RGBColor(0xA7,0x8B,0xFA)),
    ("04", "Extraction Champs",   "VLM spécialisé + Regex",     RGBColor(0x34,0xD3,0x99)),
    ("05", "Validation Métier",   "Règles domaine TOPNET",      RGBColor(0xFB,0xBF,0x24)),
    ("06", "Translitération",     "Arabe ↔ Français",           RGBColor(0xF4,0x72,0x72)),
    ("07", "Décision + BDD",      "ACCEPTÉ / REJETÉ / RÉVISER", GREEN_OK),
]
for i, (num, title, sub, col) in enumerate(steps):
    x = 0.38 + i * 1.82
    rect(sl, x, 1.55, 1.65, 4.9, fill=BLUE_LIGHT)
    rect(sl, x, 1.55, 1.65, 0.5, fill=col)
    txbox(sl, num,   x, 1.57, 1.65, 0.44, size=18, bold=True,
          color=BLUE_DARK, align=PP_ALIGN.CENTER)
    txbox(sl, title, x+0.05, 2.15, 1.55, 0.65, size=11.5, bold=True,
          color=col, align=PP_ALIGN.CENTER)
    txbox(sl, sub,   x+0.05, 2.88, 1.55, 0.8,  size=10,
          color=BLUE_LIGHT, align=PP_ALIGN.CENTER, italic=True)
    # Flèche entre étapes
    if i < 6:
        txbox(sl, "→", x+1.55, 3.7, 0.25, 0.4, size=16,
              bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# Résultat final
rect(sl, 0.38, 6.6, 12.57, 0.45, fill=BLUE_LIGHT)
txbox(sl, "✅ ACCEPTÉ  —  document authentique + champs extraits en FR/AR",
      0.5, 6.62, 6, 0.38, size=12, bold=True, color=GREEN_OK)
txbox(sl, "❌ REJETÉ  —  fraude / type inconnu / données incohérentes",
      6.6, 6.62, 6.1, 0.38, size=12, bold=True, color=RED_KO)
footer(sl)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — AVANCEMENT CE QUI EST FAIT
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
rect(sl, 0, 0, 13.33, 7.5, fill=GRAY_LIGHT)
header_bar(sl, "Avancement — Phase 1 réalisée ✅", "100% fonctionnel et validé")

done = [
    "Pipeline IA complet 7 étapes",
    "Détection fraude (ELA + VLM forensique)",
    "Classification triple couche (CLIP + VLM + OCR)",
    "OCR bilingue arabe/français (PaddleOCR)",
    "Extraction structurée par type de document",
    "Validation métier TOPNET (CIN/Passeport/Facture/Contrat)",
    "Translitération arabe ↔ français",
    "Interface Angular 21 complète (upload, résultats, historique)",
    "Authentification JWT (agent + admin)",
    "Script évaluation académique (Precision/Recall/F1/Confusion)",
    "Données synthétiques générées (passeports + contrats)",
    "Git — état stable sauvegardé",
]
todo = [
    "Migration SQLite → PostgreSQL 16",
    "Pipeline ETL + tables structurées",
    "Dashboards Power BI connectés",
    "Fine-tuning CLIP sur MIDV-500",
    "Containerisation Docker + GPU",
]

# Colonne gauche — FAIT
rect(sl, 0.35, 1.55, 7.7, 5.25, fill=WHITE)
rect(sl, 0.35, 1.55, 7.7, 0.45, fill=GREEN_OK)
txbox(sl, "✅  Ce qui est réalisé", 0.5, 1.58, 7.4, 0.38, size=13, bold=True, color=WHITE)
for i, item in enumerate(done):
    y = 2.1 + i * 0.38
    txbox(sl, "✓  " + item, 0.5, y, 7.3, 0.36, size=11, color=BLUE_DARK)

# Colonne droite — RESTE
rect(sl, 8.3, 1.55, 4.65, 5.25, fill=WHITE)
rect(sl, 8.3, 1.55, 4.65, 0.45, fill=ORANGE)
txbox(sl, "🔄  Ce qui reste (Phase 2)", 8.45, 1.58, 4.4, 0.38,
      size=13, bold=True, color=WHITE)
for i, item in enumerate(todo):
    y = 2.1 + i * 0.65
    rect(sl, 8.4, y, 4.45, 0.55, fill=GRAY_LIGHT)
    txbox(sl, item, 8.55, y+0.08, 4.2, 0.42, size=11.5, color=BLUE_DARK)
footer(sl)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — MÉTRIQUES DÉTAILLÉES
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
rect(sl, 0, 0, 13.33, 7.5, fill=GRAY_LIGHT)
header_bar(sl, "Résultats & Métriques académiques", "Évaluation sur 7 documents de test")

# Tableau résultats
headers = ["Document", "Décision", "Type détecté", "Confiance", "Fraude", "Résultat"]
rows = [
    ("cin_reel.jpg",               "ACCEPTÉ", "CIN",       "38%", "26%",  "✅ Correct"),
    ("cin_2_rectofake.png",        "REJETÉ",  "CIN",       "66%", "56%",  "✅ Correct"),
    ("cin_fake.jpeg",              "REJETÉ",  "UNKNOWN",   "0%",  "28%",  "✅ Correct"),
    ("contrat_reel.png",           "ACCEPTÉ", "CONTRAT",   "50%", "28%",  "✅ Correct"),
    ("contract_4G_5G_fake.png",    "ACCEPTÉ", "CONTRAT",   "70%", "24%",  "⚠️ Limite"),
    ("facture_real.png",           "ACCEPTÉ", "FACTURE",   "69%", "24%",  "✅ Correct"),
    ("passport_DEU_fake.png",      "REJETÉ",  "PASSEPORT", "66%", "27%",  "✅ Correct"),
]
col_w = [3.2, 1.4, 1.6, 1.3, 1.2, 1.5]
col_x = [0.3]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

# En-tête tableau
for j, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
    rect(sl, x, 1.55, w-0.05, 0.42, fill=ORANGE)
    txbox(sl, h, x+0.05, 1.58, w-0.1, 0.36, size=11, bold=True,
          color=BLUE_DARK, align=PP_ALIGN.CENTER)
# Lignes
for i, row in enumerate(rows):
    bg = BLUE_LIGHT if i % 2 == 0 else BLUE_LIGHT
    for j, (val, x, w) in enumerate(zip(row, col_x, col_w)):
        rect(sl, x, 2.05+i*0.62, w-0.05, 0.58, fill=bg)
        col = GREEN_OK if "✅" in val else (ORANGE if "⚠️" in val
              else (RED_KO if val == "REJETÉ" else WHITE))
        txbox(sl, val, x+0.05, 2.1+i*0.62, w-0.1, 0.5, size=10.5,
              color=col, align=PP_ALIGN.CENTER)

# Récap métriques
rect(sl, 0.3, 6.45, 12.7, 0.55, fill=BLUE_LIGHT)
summary = [
    ("Accuracy : 85.7%", GREEN_OK), ("Precision : 100%", GREEN_OK),
    ("Recall : 75%", ORANGE), ("F1 : 85.7%", GREEN_OK),
    ("Faux positifs : 0%", GREEN_OK),
]
for i, (txt, col) in enumerate(summary):
    txbox(sl, txt, 0.5+i*2.5, 6.5, 2.4, 0.42, size=12, bold=True,
          color=col, align=PP_ALIGN.CENTER)
footer(sl)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — DIFFICULTÉS
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
rect(sl, 0, 0, 13.33, 7.5, fill=GRAY_LIGHT)
header_bar(sl, "Difficultés rencontrées & Solutions")

difficulties = [
    ("⏱️  Performance CPU",
     "VLM Qwen2.5vl : 60-180s/image sans GPU",
     "Fast-path : skip VLM si OCR+CLIP confiants → -30-40% de temps"),
    ("🔤  OCR arabe bilingues",
     "CIN : texte arabe + français simultané difficile à extraire",
     "PaddleOCR + EasyOCR en parallèle + fusion des résultats"),
    ("📄  Faux PDF numériques",
     "Un PDF bien forgé a le même profil ELA qu'un vrai PDF",
     "Limite documentée — nécessite validation BDD TOPNET (phase 2)"),
    ("🔄  Angular change detection",
     "Timer et résultats ne s'affichaient pas (hors NgZone)",
     "NgZone.run() + ChangeDetectorRef.detectChanges() appliqués"),
]
for i, (title, prob, sol) in enumerate(difficulties):
    row, col = divmod(i, 2)
    x = 0.35 + col * 6.5
    y = 1.55 + row * 2.6
    rect(sl, x, y, 6.2, 2.35, fill=WHITE)
    rect(sl, x, y, 6.2, 0.45, fill=BLUE_DARK)
    txbox(sl, title, x+0.1, y+0.05, 6.0, 0.38, size=13, bold=True, color=ORANGE)
    txbox(sl, "❌ " + prob, x+0.1, y+0.55, 6.0, 0.55, size=11.5, color=RED_KO, italic=True)
    rect(sl, x+0.1, y+1.15, 0.05, 0.85, fill=GREEN_OK)
    txbox(sl, "✅ " + sol, x+0.25, y+1.15, 5.85, 0.85, size=11.5, color=BLUE_DARK)
footer(sl)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CONCLUSION
# ═════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
rect(sl, 0, 0, 13.33, 2.2, fill=BLUE_DARK)
rect(sl, 0, 2.2, 13.33, 0.07, fill=ORANGE)
rect(sl, 0, 0, 0.35, 7.5, fill=ORANGE)

txbox(sl, "Conclusion & Perspectives", 0.6, 0.25, 12.5, 0.75,
      size=30, bold=True, color=WHITE)
txbox(sl, "Phase 1 complète — Pipeline IA opérationnel", 0.6, 0.92, 12, 0.55,
      size=18, color=ORANGE)

# Acquis
rect(sl, 0.5, 1.65, 5.8, 1.8, fill=BLUE_LIGHT)
rect(sl, 0.5, 1.65, 5.8, 0.42, fill=GREEN_OK)
txbox(sl, "✅  Objectifs Phase 1 atteints", 0.65, 1.68, 5.5, 0.36,
      size=13, bold=True, color=BLUE_DARK)
multiline_box(sl, [
    ("Pipeline IA 7 étapes fonctionnel", False, WHITE),
    ("Précision 100% — 0% faux positifs", True,  GREEN_OK),
    ("Interface web complète et opérationnelle", False, WHITE),
], 0.65, 2.17, 5.5, 1.2, size=12)

# Compétences
rect(sl, 6.8, 1.65, 6.1, 1.8, fill=BLUE_LIGHT)
rect(sl, 6.8, 1.65, 6.1, 0.42, fill=RGBColor(0x38,0xBD,0xF8))
txbox(sl, "📚  Compétences acquises", 6.95, 1.68, 5.8, 0.36,
      size=13, bold=True, color=BLUE_DARK)
multiline_box(sl, [
    ("Computer Vision : CLIP, VLM, ELA, OCR", False, WHITE),
    ("MLOps : métriques, tuning, évaluation", False, WHITE),
    ("Full-stack : FastAPI + Angular + BDD", False, WHITE),
], 6.95, 2.17, 5.8, 1.2, size=12)

# Perspectives
rect(sl, 0.5, 3.65, 12.4, 2.55, fill=BLUE_LIGHT)
rect(sl, 0.5, 3.65, 12.4, 0.42, fill=ORANGE)
txbox(sl, "🚀  Perspectives Phase 2", 0.65, 3.68, 12.1, 0.36,
      size=13, bold=True, color=BLUE_DARK)
persp = [
    ("PostgreSQL + ETL", "Données structurées pour analytique"),
    ("Power BI Dashboards", "KPIs temps réel pour management"),
    ("Fine-tuning CLIP", "MIDV-500 → meilleure précision passeport"),
    ("Docker + GPU", "Temps 121s → ~8s en production"),
]
for i, (title, desc) in enumerate(persp):
    x = 0.65 + i * 3.05
    txbox(sl, title, x, 4.18, 2.9, 0.38, size=12, bold=True, color=ORANGE)
    txbox(sl, desc,  x, 4.58, 2.9, 0.55, size=11, color=BLUE_LIGHT)

txbox(sl, "Merci pour votre attention — Questions ?",
      0.6, 6.45, 12, 0.55, size=20, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)
footer(sl)

# ── Sauvegarde ────────────────────────────────────────────────────────────────
output = "TOPNET_OCR_Presentation_PFE.pptx"
prs.save(output)
print(f"✅ Présentation générée : {output}")
print(f"   10 slides | Design TOPNET bleu #1e3a5f + orange #ff9500")
print(f"   Ouvrir dans PowerPoint ou importer dans Canva/Google Slides")
